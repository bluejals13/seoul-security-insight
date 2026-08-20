import os
import sys

def create_presentation():
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.enum.shapes import MSO_SHAPE
    except ImportError:
        print("[NOTICE] 'python-pptx' module is not installed.")
        print("To generate presentation/seoul-security-insight.pptx, run:")
        print("    uv add python-pptx  (or pip install python-pptx)")
        print("    python scripts/generate_pptx.py")
        return False

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    NAVY = RGBColor(15, 23, 42)
    WHITE = RGBColor(255, 255, 255)
    BLUE = RGBColor(59, 130, 246)
    PURPLE = RGBColor(168, 85, 247)
    SLATE_GRAY = RGBColor(148, 163, 184)
    ROSE = RGBColor(244, 63, 94)

    def add_background(slide, color=NAVY):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = color
        bg.line.fill.background()
        return bg

    def add_header(slide, title_text, category_text="SEOUL SECURITY INSIGHT"):
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.8))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p0 = tf.paragraphs[0]
        p0.text = category_text.upper()
        p0.font.size = Pt(10)
        p0.font.bold = True
        p0.font.color.rgb = BLUE
        
        p1 = tf.add_paragraph()
        p1.text = title_text
        p1.font.size = Pt(24)
        p1.font.bold = True
        p1.font.color.rgb = WHITE

    # SLIDE 1: Title
    s1 = prs.slides.add_slide(blank_layout)
    add_background(s1, NAVY)
    tb1 = s1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(3.8))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    
    p = tf1.paragraphs[0]
    p.text = "서울은 어디가 더 안전한가?"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf1.add_paragraph()
    p2.text = "데이터 기반 팩트, 통계 분석, 그리고 객관적 한계의 구별"
    p2.font.size = Pt(18)
    p2.font.color.rgb = SLATE_GRAY
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(15)

    p3 = tf1.add_paragraph()
    p3.text = "[FACT] 2024년 5대 범죄 80,819건 전수 | [INTERPRETATION] 인구 보정 | [LIMITATION] 인과관계 제외"
    p3.font.size = Pt(12)
    p3.font.color.rgb = BLUE
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(30)

    # SLIDE 2: Problem & Framework
    s2 = prs.slides.add_slide(blank_layout)
    add_background(s2, NAVY)
    add_header(s2, "범죄 발생 건수가 많다고 정말 더 위험한 지역일까?", "02. PROBLEM STATEMENT")
    
    tb2 = s2.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    
    p = tf2.paragraphs[0]
    p.text = "[FACT] 데이터 사실"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BLUE
    
    p_fact = tf2.add_paragraph()
    p_fact.text = "• 강남구는 2024년 5대 범죄 절대 발생건수 1위 (6,107건)입니다.\n"
    p_fact.font.size = Pt(13)
    p_fact.font.color.rgb = WHITE

    p_interp_t = tf2.add_paragraph()
    p_interp_t.text = "[INTERPRETATION] 합리적 해석"
    p_interp_t.font.size = Pt(16)
    p_interp_t.font.bold = True
    p_interp_t.font.color.rgb = PURPLE
    p_interp_t.space_before = Pt(10)
    
    p_interp = tf2.add_paragraph()
    p_interp.text = "• 절대 건수만으로 치안 위험도를 단정하는 것은 지역 인구 규모를 고려하지 않은 단순 합산의 한계가 존재합니다.\n"
    p_interp.font.size = Pt(13)
    p_interp.font.color.rgb = WHITE

    p_lim_t = tf2.add_paragraph()
    p_lim_t.text = "[LIMITATION] 데이터 한계"
    p_lim_t.font.size = Pt(16)
    p_lim_t.font.bold = True
    p_lim_t.font.color.rgb = ROSE
    p_lim_t.space_before = Pt(10)
    
    p_lim = tf2.add_paragraph()
    p_lim.text = "• 현재 상주인구 데이터만으로는 실질 치안 위험도를 완벽히 판단할 수 없으며 유동인구 등 추가 데이터 검증이 필요합니다."
    p_lim.font.size = Pt(13)
    p_lim.font.color.rgb = WHITE

    # SLIDE 3: Key Insight 01
    s3 = prs.slides.add_slide(blank_layout)
    add_background(s3, NAVY)
    add_header(s3, "KEY INSIGHT 01: 절대 1위(강남) vs 인구 1만명당 1위(중구)", "05. KEY INSIGHT 01")
    
    tb3 = s3.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf3 = tb3.text_frame
    tf3.word_wrap = True
    
    p = tf3.paragraphs[0]
    p.text = "[FACT] 강남구 절대 1위 (6,107건) vs 중구 인구 1만명당 1위 (454.65건)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BLUE
    
    p_stat = tf3.add_paragraph()
    p_stat.text = "• 통계 결과: Pearson r = 0.186, Spearman rho = 0.418 (순위 변동 존재)"
    p_stat.font.size = Pt(13)
    p_stat.font.color.rgb = WHITE

    p_interp3 = tf3.add_paragraph()
    p_interp3.text = "\n[INTERPRETATION] 중구·종로구는 상주인구가 상대적으로 적기 때문에, 상주인구를 분모로 사용한 범죄율이 높게 계산됩니다. 실제 생활인구·유동인구를 반영할 경우 결과가 달라질 가능성이 있습니다."
    p_interp3.font.size = Pt(13)
    p_interp3.font.color.rgb = PURPLE

    p_lim3 = tf3.add_paragraph()
    p_lim3.text = "\n[LIMITATION] 본 분석에는 유동인구·생활인구 데이터가 포함되어 있지 않아, 높은 범죄율의 원인을 유동인구 때문이라고 단정할 수 없습니다."
    p_lim3.font.size = Pt(13)
    p_lim3.font.color.rgb = ROSE

    # SLIDE 4: Key Insight 02 (Exact Crime Breakdown)
    s4 = prs.slides.add_slide(blank_layout)
    add_background(s4, NAVY)
    add_header(s4, "KEY INSIGHT 02: 5대 범죄의 92.9%는 폭력과 절도", "06. KEY INSIGHT 02")
    
    tb4 = s4.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf4 = tb4.text_frame
    tf4.word_wrap = True
    
    p = tf4.paragraphs[0]
    p.text = "[FACT] 2024년 5대 범죄 총 80,819건의 정밀 범죄 유형별 수치"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BLUE
    
    p_break = tf4.add_paragraph()
    p_break.text = "• 폭력: 39,554건 (48.9%)\n• 절도: 35,506건 (43.9%)\n• 성범죄: 5,523건 (6.8%)\n• 살인: 149건 (0.18%)\n• 강도: 87건 (0.11%)\n• 폭력 + 절도 합계: 75,060건 (92.9%)"
    p_break.font.size = Pt(13)
    p_break.font.color.rgb = WHITE

    p_interp4 = tf4.add_paragraph()
    p_interp4.text = "\n[INTERPRETATION] 자치구별로 범죄 유형의 구성비가 서로 다르게 나타나는 패턴이 관찰됩니다."
    p_interp4.font.size = Pt(13)
    p_interp4.font.color.rgb = PURPLE

    p_interp4 = tf4.add_paragraph()
    p_interp4.text = "\n[LIMITATION] 본 데이터만으로 이러한 차이가 상권, 주거형태 등 특정 지역 특성에 의해 발생했다고 판단할 수 없습니다."
    p_interp4.font.size = Pt(13)
    p_interp4.font.color.rgb = ROSE

    # SLIDE 5: Key Insight 03 (Streetlights)
    s5 = prs.slides.add_slide(blank_layout)
    add_background(s5, NAVY)
    add_header(s5, "KEY INSIGHT 03: 가로등 수 vs 범죄 발생량 분석", "07. KEY INSIGHT 03")
    
    tb5 = s5.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf5 = tb5.text_frame
    tf5.word_wrap = True
    
    p = tf5.paragraphs[0]
    p.text = "[FACT] 가로등 수 vs 범죄 건수: Pearson r = 0.561, Spearman rho = 0.523"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BLUE
    
    p_stat5 = tf5.add_paragraph()
    p_stat5.text = "• 인구 1만 명당 가로등 수 vs 1만 명당 범죄율: Pearson r = 0.120, Spearman rho = 0.287 (상관성 없음/낮음)"
    p_stat5.font.size = Pt(13)
    p_stat5.font.color.rgb = WHITE

    p_interp5 = tf5.add_paragraph()
    p_interp5.text = "\n[INTERPRETATION] 자치구별 가로등 수와 범죄 발생건수의 절대량 사이에는 양의 상관관계가 관찰됩니다. 그러나 인구를 보정하면 상관계수가 낮아집니다."
    p_interp5.font.size = Pt(13)
    p_interp5.font.color.rgb = PURPLE

    p_lim5 = tf5.add_paragraph()
    p_lim5.text = "\n[LIMITATION] 상관관계는 인과관계를 의미하지 않으며, 가로등이 범죄를 증가시키거나 감소시킨다는 결론을 내릴 수 없습니다."
    p_lim5.font.size = Pt(13)
    p_lim5.font.color.rgb = ROSE

    os.makedirs("presentation", exist_ok=True)
    output_path = "presentation/seoul-security-insight.pptx"
    prs.save(output_path)
    print(f"Factchecked PPTX generated successfully at: {output_path}")
    return True

if __name__ == "__main__":
    create_presentation()
